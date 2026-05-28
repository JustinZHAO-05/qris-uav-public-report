from __future__ import annotations

import json
import re
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures_v2"
ASSETS = ROOT / "assets" / "gpt_image2"
OUT = ROOT / "outputs" / "ppt_v2"
REPORTS = ROOT / "outputs" / "reports_v2"
TABLES = ROOT / "outputs" / "tables_v2"
FORMULAS = ROOT / "outputs" / "formulas_v2"
FINAL = OUT / "Q-RIS-UAV-public-academic-report.pptx"

W, H = 13.333, 7.5
TOTAL_SLIDES = 77
CURRENT_SECTION = ""

INK = RGBColor(22, 32, 42)
MUTED = RGBColor(86, 96, 106)
PAPER = RGBColor(249, 247, 241)
PANEL = RGBColor(247, 249, 250)
DEEP = RGBColor(29, 43, 57)
TEAL = RGBColor(42, 157, 143)
CORAL = RGBColor(231, 111, 81)
GOLD = RGBColor(233, 196, 106)
BLUE = RGBColor(69, 123, 157)
GRAY = RGBColor(141, 153, 174)
WHITE = RGBColor(255, 255, 255)
PURPLE = RGBColor(102, 80, 164)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"
FONT_MONO = "Consolas"


def set_run(run, size=16, color=INK, bold=False, font=FONT_CN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False, align="left", font=FONT_CN, leading=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = leading
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color, bold=bold, font=font)
    return box


def add_bg(slide, color=PAPER):
    r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()


def add_panel(slide, x, y, w, h, fill=PANEL, line=RGBColor(218, 224, 228)):
    r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.color.rgb = line
    r.line.width = Pt(0.7)
    return r


def add_picture(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    img_aspect = iw / ih
    box_aspect = w / h
    if img_aspect >= box_aspect:
        pic_w = w
        pic_h = w / img_aspect
    else:
        pic_h = h
        pic_w = h * img_aspect
    px = x + (w - pic_w) / 2
    py = y + (h - pic_h) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))


def set_section(section: str) -> None:
    global CURRENT_SECTION
    CURRENT_SECTION = section


def add_title(slide, kicker, title, subtitle=None, dark=False):
    color = WHITE if dark else INK
    muted = RGBColor(210, 218, 224) if dark else MUTED
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.42), Inches(0.34), Inches(0.035)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = TEAL
    slide.shapes[-1].line.fill.background()
    add_text(slide, kicker.upper(), 0.98, 0.31, 2.4, 0.26, 8.5, muted, True, font=FONT_EN)
    if CURRENT_SECTION:
        add_text(slide, CURRENT_SECTION, 8.10, 0.31, 4.30, 0.26, 8.5, muted, True, align="right")
    add_text(slide, title, 0.55, 0.64, 11.6, 0.55, 26, color, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.16, 10.8, 0.36, 12.5, muted)


def add_footer(slide, page, total=TOTAL_SLIDES, text="Q-RIS-UAV-AO · quantum electromagnetic SAGSI system"):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.02), Inches(12.15), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(214, 219, 223)
    line.line.fill.background()
    add_text(slide, text, 0.55, 7.08, 7.0, 0.20, 8, GRAY, font=FONT_EN)
    add_text(slide, f"{page:02d}/{total}", 12.0, 7.08, 0.70, 0.20, 8, GRAY, align="right", font=FONT_EN)


def add_bullets(slide, items, x, y, w, h, size=15, accent=TEAL):
    row_h = h / max(len(items), 1)
    for i, item in enumerate(items):
        yy = y + i * row_h
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.08), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()
        add_text(slide, item, x + 0.22, yy, w - 0.22, row_h * 0.92, size=size, color=INK)


def add_metric(slide, value, label, x, y, w, accent=TEAL, dark=False):
    add_text(slide, value, x, y, w, 0.42, 24, WHITE if dark else accent, True, font=FONT_EN)
    add_text(slide, label, x, y + 0.42, w, 0.34, 10.5, RGBColor(205, 214, 220) if dark else MUTED)


def _color_to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def formula_png(name: str, lines: list[str], color=WHITE, fontsize=23) -> Path:
    FORMULAS.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[^0-9A-Za-z_\-]+", "_", name).strip("_").lower()
    path = FORMULAS / f"{safe}.png"
    height = 0.75 + 0.80 * len(lines)
    fig = plt.figure(figsize=(8.3, height), dpi=260)
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    y = 0.94
    step = 0.88 / max(len(lines), 1)
    for line in lines:
        ax.text(0.02, y, f"${line}$", fontsize=fontsize, color=_color_to_hex(color), va="top")
        y -= step
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return path


def add_formula(slide, title, formula, x, y, w, h, accent=TEAL, formula_id=None):
    add_panel(slide, x, y, w, h, DEEP, DEEP)
    add_text(slide, title, x + 0.25, y + 0.22, w - 0.50, 0.34, 15, accent, True)
    if isinstance(formula, str):
        lines = [line.strip() for line in formula.split("\n") if line.strip()]
    else:
        lines = list(formula)
    image = formula_png(formula_id or title, lines, WHITE, fontsize=30 if len(lines) <= 4 else 25)
    add_picture(slide, image, x + 0.25, y + 0.72, w - 0.50, h - 1.02)


def add_math_card(slide, formula, x, y, w, h, accent=TEAL, formula_id=None, fontsize=28):
    add_panel(slide, x, y, w, h, DEEP, DEEP)
    if isinstance(formula, str):
        lines = [line.strip() for line in formula.split("\n") if line.strip()]
    else:
        lines = list(formula)
    image = formula_png(formula_id or "_".join(lines)[:24], lines, WHITE, fontsize=fontsize)
    add_picture(slide, image, x + 0.18, y + 0.16, w - 0.36, h - 0.32)


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=10.2):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, ratio in enumerate(col_widths):
            table.columns[i].width = Inches(w * ratio)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DEEP if r == 0 else RGBColor(252, 252, 249)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, font_size, WHITE if r == 0 else INK, r == 0)
    return table_shape


def add_explanation_panel(slide, heading, body, x, y, w, h, accent=TEAL, dark=False, body_size=11.2):
    fill = RGBColor(252, 252, 249) if not dark else RGBColor(42, 58, 73)
    line = RGBColor(218, 224, 228) if not dark else RGBColor(78, 95, 108)
    text_color = INK if not dark else RGBColor(218, 226, 231)
    add_panel(slide, x, y, w, h, fill, line)
    add_text(slide, heading, x + 0.22, y + 0.18, w - 0.44, 0.30, 13.5, accent, True)
    add_text(slide, body, x + 0.22, y + 0.58, w - 0.44, h - 0.75, body_size, text_color, leading=1.16)


def image_slide(prs, page, kicker, title, image, subtitle=None, note="", dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP if dark else PAPER)
    add_title(slide, kicker, title, dark=dark)
    add_picture(slide, image, 0.60, 1.40, 8.70, 5.35)
    panel_fill = RGBColor(252, 252, 249) if not dark else RGBColor(42, 58, 73)
    panel_line = RGBColor(218, 224, 228) if not dark else RGBColor(78, 95, 108)
    add_panel(slide, 9.52, 1.40, 2.95, 5.35, panel_fill, panel_line)
    label = "图表解读" if str(kicker).upper().startswith("RESULT") or Path(image).parent == FIG else "机制说明"
    label_color = TEAL if not dark else GOLD
    text_color = INK if not dark else RGBColor(218, 226, 231)
    add_text(slide, label, 9.72, 1.68, 1.65, 0.30, 15.0, label_color, True)
    body = subtitle or "该页用于说明系统模块、实验对象与后续建模之间的关系。"
    add_text(slide, body, 9.72, 2.15, 2.38, 4.05, 11.0, text_color, leading=1.14)
    add_footer(slide, page, text=note or "Q-RIS-UAV-AO · quantum electromagnetic SAGSI system")
    return slide


def transition_slide(prs, page, section_no, section_title, lead, bullets):
    set_section(f"{section_no}、{section_title}")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP)
    add_text(slide, f"{section_no}", 0.70, 0.78, 1.05, 0.55, 26, TEAL, True, font=FONT_EN)
    add_text(slide, section_title, 1.82, 0.72, 8.8, 0.60, 28, WHITE, True)
    add_text(slide, lead, 1.85, 1.52, 9.9, 0.55, 15, RGBColor(218, 226, 231))
    for i, item in enumerate(bullets):
        y = 2.55 + i * 0.72
        add_text(slide, f"{i + 1}", 1.88, y, 0.32, 0.26, 12, GOLD, True, font=FONT_EN)
        add_text(slide, item, 2.35, y - 0.02, 8.80, 0.34, 16, WHITE)
    add_footer(slide, page, text="Section transition")
    return slide


def analysis_slide(prs, page, kicker, title, paragraphs, takeaway, section_label=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP if dark else PAPER)
    if section_label:
        set_section(section_label)
    add_title(slide, kicker, title, dark=dark)
    x0, y0 = 0.85, 1.55
    w, h = 11.65, 4.75
    fill = RGBColor(252, 252, 249) if not dark else RGBColor(42, 58, 73)
    line = RGBColor(218, 224, 228) if not dark else RGBColor(78, 95, 108)
    add_panel(slide, x0, y0, w, h, fill, line)
    text_color = INK if not dark else RGBColor(218, 226, 231)
    y = y0 + 0.28
    for i, p in enumerate(paragraphs):
        add_text(slide, p, x0 + 0.35, y, w - 0.70, 0.92, 12.4, text_color, leading=1.14)
        y += 1.28
    add_text(slide, takeaway, 0.95, 6.28, 11.25, 0.52, 14.0, TEAL if not dark else GOLD, True, leading=1.12)
    add_footer(slide, page, text="Academic narrative bridge")
    return slide


def build() -> Path:
    OUT.mkdir(parents=True, exist_ok=True)
    REPORTS.mkdir(parents=True, exist_ok=True)
    FORMULAS.mkdir(parents=True, exist_ok=True)
    with (TABLES / "summary_metrics_v2.json").open("r", encoding="utf-8") as f:
        metrics = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    notes = []
    page = 1

    def new_slide(bg=PAPER):
        s = prs.slides.add_slide(blank)
        add_bg(s, bg)
        return s

    # 1
    set_section("封面")
    s = new_slide(DEEP)
    add_picture(s, ASSETS / "disaster_scenario_gpt.png", 6.9, 0.68, 5.85, 4.55)
    add_text(s, CURRENT_SECTION, 11.10, 0.35, 1.20, 0.24, 8.5, RGBColor(214, 224, 229), True, align="right")
    add_text(s, "量子密钥保护的 RIS-UAV 远海风电与海缆监测", 0.65, 0.82, 6.2, 0.48, 19, WHITE, True)
    add_text(s, "空天地海一体化系统：\n能效-安全联合优化设计", 0.65, 1.42, 6.2, 1.28, 30, WHITE, True)
    add_text(s, "Optimization Theory for Machine Learning - Final Project\nYanzhe Zhao\nTianjin University, Future Technology College", 0.68, 2.95, 5.75, 0.98, 13, RGBColor(214, 224, 229), leading=1.12)
    add_metric(s, f"{metrics['legacy_proposed_gain_vs_random']:.2f}x", "安全能效提升", 0.72, 5.35, 1.8, TEAL, True)
    add_metric(s, f"{metrics['dinkelbach_final_eta']:.2f}", "最终安全能效", 2.80, 5.35, 2.0, GOLD, True)
    add_metric(s, f"{metrics['magnetic_roc_auc']:.3f}", "磁异常检测AUC", 5.05, 5.35, 2.1, CORAL, True)
    add_footer(s, page, text="DeepOWT · IBTrACS · WMM2025 · MATLAB/Python RIS-UAV simulation")
    notes.append("封面：项目一句话定位和三个关键结果。")
    page += 1

    # 2
    set_section("汇报目录")
    s = new_slide()
    add_title(s, "CONTENTS", "汇报目录")
    agenda = [
        ("一、灾后远海监测场景", "风电场、海缆、台风轨迹和真实数据构造"),
        ("二、空天地海系统架构", "卫星、UAV、岸基、RIS、AUV与密钥管理闭环"),
        ("三、优化建模", "信道、QKD、磁异常优先级、能耗与完整分式目标"),
        ("四、Q-RIS-UAV-AO算法", "Dinkelbach、调度、CVXPY、RIS相位、UAV轨迹SCA"),
        ("五、仿真实验与分析", "RIS/MATLAB波束、轨迹、收敛、密钥、鲁棒性与消融"),
        ("六、创新点与工程边界", "量子感知、量子安全与可重构传播的联合优化"),
    ]
    for i, (t, b) in enumerate(agenda):
        y = 1.42 + i * 0.80
        add_text(s, f"{i + 1:02d}", 0.92, y, 0.55, 0.34, 17, TEAL, True, font=FONT_EN)
        add_text(s, t, 1.58, y - 0.02, 3.95, 0.32, 16, INK, True)
        add_text(s, b, 5.65, y - 0.02, 6.35, 0.34, 14.5, MUTED)
    add_footer(s, page, text="Academic presentation outline")
    page += 1

    transition_slide(
        prs,
        page,
        "一",
        "灾后远海监测场景",
        "先说明问题从哪里来：灾害、海缆风险、远海通信和真实数据如何进入同一场景。",
        ["远海风电和海缆是典型的高价值、低可达基础设施。", "台风后常规通信与人工巡检能力下降，无人机（Unmanned Aerial Vehicle, UAV）临时组网成为关键手段。", "量子磁力计探测提供异常线索，后续优化围绕异常优先级展开。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "INTRODUCTION",
        "问题引入：灾后远海监测为何形成优化问题",
        [
            "远海风电场和海底电缆位于人工巡检半径之外，台风后同时面临设备状态未知、通信链路退化和安全告警时效性下降的问题。此时任务目标从单纯收集更多数据，转向在有限能量和有限密钥下优先回传与风险最高区域相关的数据。",
            "本项目把灾后巡检抽象为一个联合优化问题：量子磁探测给出风险权重，RIS和UAV决定远海链路质量，QKD密钥决定哪些数据能够被安全计入收益。三个模块共同决定安全数据能效。",
            "因此，后续建模会依次回答三个问题：哪些节点重要、如何把它们连上、在能量和密钥受限时如何分配资源。这个叙事顺序对应后续的系统架构、数学模型和算法分解。",
        ],
        "本章作用：把工程场景转化为“风险权重、通信链路、安全密钥、能耗约束”四类数学对象。",
    )
    page += 1

    # 4
    s = new_slide()
    add_title(s, "TERMS", "核心术语与简写表")
    rows = [
        ["简写", "全称", "本项目中的含义"],
        ["RIS", "智能反射超表面 Reconfigurable Intelligent Surface", "通过离散相位编码调控反射链路"],
        ["UAV", "无人机 Unmanned Aerial Vehicle", "灾后临时空中中继与巡检平台"],
        ["QKD", "量子密钥分发 Quantum Key Distribution", "为关键告警与控制数据提供密钥预算"],
        ["SAGSI", "空天地海一体化网络 Space-Air-Ground-Sea Integrated Network", "卫星、UAV、岸基、海上/水下节点协同"],
        ["AUV", "自主水下航行器 Autonomous Underwater Vehicle", "搭载量子磁力计进行海缆异常探测"],
    ]
    add_table(s, rows, 0.70, 1.42, 11.90, 4.55, col_widths=[0.13, 0.43, 0.44], font_size=10.8)
    add_text(s, "本页先给出跨学科术语的中文含义与英文全称。后续章节在第一次进入数学模型时再定义变量符号，使系统概念、工程对象和优化变量保持一致。", 0.95, 6.03, 11.0, 0.55, 14.5, MUTED)
    add_footer(s, page)
    page += 1

    # 5
    s = new_slide()
    add_title(s, "PROJECT SPINE", "研究主线：量子电磁感知、安全通信与能效优化闭环")
    add_panel(s, 0.85, 1.55, 3.65, 3.95, RGBColor(239, 249, 247))
    add_text(s, "感知", 1.12, 1.88, 1.3, 0.32, 20, TEAL, True)
    add_text(s, "自主水下航行器（AUV）搭载量子磁力计，对海缆附近磁扰动进行测量。观测磁场先扣除WMM2025背景，再映射为异常分数和任务权重，决定哪些节点的数据更应优先回传。", 1.12, 2.22, 3.05, 1.95, 12.2, INK, leading=1.12)
    add_text(s, "建模入口：异常分数 𝐴ₖ[𝑛] 与权重 𝑤ₖ[𝑛]", 1.12, 4.55, 3.05, 0.44, 11.6, TEAL, True)
    add_panel(s, 4.85, 1.55, 3.65, 3.95, RGBColor(255, 247, 238))
    add_text(s, "通信", 5.12, 1.88, 1.3, 0.32, 20, CORAL, True)
    add_text(s, "智能反射超表面（RIS）通过离散相位改变级联信道，使远海节点、RIS与UAV之间的反射链路相干叠加。UAV同时通过轨迹、功率和带宽分配完成临时回传。", 5.12, 2.22, 3.05, 1.95, 12.2, INK, leading=1.12)
    add_text(s, "建模入口：速率 𝑅ₖ[𝑛]、相位矩阵 𝜣[𝑛]", 5.12, 4.55, 3.05, 0.44, 11.6, CORAL, True)
    add_panel(s, 8.85, 1.55, 3.65, 3.95, RGBColor(242, 246, 251))
    add_text(s, "安全", 9.12, 1.88, 1.3, 0.32, 20, BLUE, True)
    add_text(s, "卫星量子密钥分发（QKD）提供有限密钥预算。关键告警和控制指令只有在获得足够密钥时才能计入安全数据量，因此密钥供给会与无线容量共同形成瓶颈", 9.12, 2.22, 3.05, 1.95, 12.2, INK, leading=1.12)
    add_text(s, "建模入口：密钥分配 κₖ[𝑛] 与安全量 𝑠ₖ[𝑛]", 9.12, 4.55, 3.05, 0.44, 11.6, BLUE, True)
    add_text(s, "核心闭环：磁异常 → 任务权重 → RIS/UAV/资源联合优化 → 安全数据能效最大化", 1.05, 6.10, 11.0, 0.40, 20, INK, True)
    add_footer(s, page)
    notes.append("总论点：所有技术最后落到优化变量、目标和约束中。")
    page += 1

    image_slide(prs, page, "BACKGROUND", "灾后远海风电与海缆监测场景", ASSETS / "disaster_scenario_gpt.png", "该机制图描述台风后远海风电场、海缆和通信链路同时受扰的应急场景。岸基站提供边缘计算和密钥管理，UAV承担临时中继，AUV负责水下磁异常探测，RIS用于增强远海无线链路。该场景引出后续优化目标：在能量、带宽和密钥均受限的条件下，优先回传与海缆风险相关的安全数据。", "Scenario diagram based on DeepOWT and IBTrACS")
    notes.append("背景：灾后场景和关键问题。")
    page += 1

    # 4
    s = new_slide()
    add_title(s, "MOTIVATION", "空天地海协同监测的必要性")
    rows = [
        ["层级", "关键能力", "本项目承担的任务"],
        ["空：卫星", "SAR遥感 + QKD", "灾后粗定位、密钥分发"],
        ["天：UAV", "移动中继 + 巡检", "临时通信、轨迹优化"],
        ["地：岸基/电网", "边缘计算 + KMS", "优化求解、安全策略下发"],
        ["海：平台/AUV/RIS", "水下感知 + 可控传播", "磁异常检测、链路增强"],
    ]
    add_table(s, rows, 0.85, 1.55, 11.55, 3.80, col_widths=[0.18, 0.30, 0.52], font_size=13)
    add_text(s, "设计边界：卫星侧提供遥感与密钥，UAV侧提供临时空中链路，岸基侧完成优化求解，海上/水下侧提供感知与可重构传播。四层功能共同决定后续变量、目标函数和约束。", 1.05, 5.88, 11.0, 0.65, 14.5, CORAL, True, leading=1.12)
    add_footer(s, page)
    notes.append("说明空天地海四层各自承担不可替代的工程任务。")
    page += 1

    s = new_slide()
    add_title(s, "DATA LOGIC", "真实数据、物理背景与可控仿真的组合方式")
    add_panel(s, 0.85, 1.50, 3.55, 4.55, RGBColor(239, 249, 247))
    add_text(s, "真实背景", 1.12, 1.82, 1.6, 0.32, 17, TEAL, True)
    add_bullets(s, ["DeepOWT 给出海上风电基础设施空间分布", "NOAA IBTrACS 给出台风轨迹背景", "Natural Earth 提供海岸线底图"], 1.10, 2.30, 2.95, 2.50, size=12.0, accent=TEAL)
    add_panel(s, 4.85, 1.50, 3.55, 4.55, RGBColor(255, 247, 238))
    add_text(s, "感知背景", 5.12, 1.82, 1.6, 0.32, 17, CORAL, True)
    add_bullets(s, ["WMM2025 给出地磁背景尺度", "海缆异常用偶极扰动构造", "异常分数转化为调度权重"], 5.10, 2.30, 2.95, 2.50, size=12.0, accent=CORAL)
    add_panel(s, 8.85, 1.50, 3.55, 4.55, RGBColor(242, 246, 251))
    add_text(s, "可控仿真", 9.12, 1.82, 1.6, 0.32, 17, BLUE, True)
    add_bullets(s, ["通信信道和RIS采用可复现实验仿真", "优化算法固定随机种子", "每个结果图都可由脚本再生成"], 9.10, 2.30, 2.95, 2.50, size=12.0, accent=BLUE)
    add_text(s, "方法说明：真实公开数据用于约束场景空间位置与灾害背景，物理模型用于给出磁场、链路和能耗的量纲关系；可控仿真用于比较算法变量变化带来的因果影响。这样既避免把全部结果建立在纯合成场景上，也避免真实数据缺少通信控制变量的问题。", 0.95, 6.08, 11.15, 0.70, 12.4, MUTED, leading=1.12)
    add_footer(s, page)
    page += 1

    image_slide(prs, page, "DATA", "数据来源矩阵", FIG / "data_source_matrix.png", "该表将公开数据与仿真实验逐一对应。DeepOWT用于选择海上风电节点，IBTrACS用于提供台风轨迹背景，Natural Earth用于海岸线底图，WMM2025用于地磁背景尺度。通信信道与RIS部分采用固定随机种子的可控仿真，以便隔离算法贡献并保证结果可复现。", "DeepOWT · IBTrACS · WMM2025 · Natural Earth · Sentinel-1")
    notes.append("数据源页：真实数据和合成仿真的分工。")
    page += 1

    image_slide(prs, page, "MAP", "真实台风轨迹、风电节点与海缆风险区域叠加图", FIG / "typhoon_deepowt_context.png", "该地图由NOAA IBTrACS台风轨迹、Natural Earth海岸线和DeepOWT海上风电坐标共同构成。图中节点来自真实地理背景下筛选出的远海基础设施；岸基站、RIS平台和海缆异常点进一步定义了通信、感知和优化的空间关系。", "NOAA IBTrACS + Natural Earth coastline + DeepOWT")
    notes.append("地图页：新增IBTrACS和Natural Earth。")
    page += 1

    transition_slide(
        prs,
        page,
        "二",
        "空天地海系统架构",
        "把灾后任务拆给不同层级节点：卫星负责态势和密钥，UAV负责临时链路，海上/水下节点负责感知与可重构传播。",
        ["空：卫星遥感和QKD密钥供给。", "天：UAV移动中继与轨迹优化。", "海：AUV磁探测、海上RIS平台和风电节点数据回传。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "SYSTEM INTRO",
        "系统架构导言：从节点分工到优化变量",
        [
            "空天地海一体化系统的核心在于把不同空间层级的能力分配给明确的物理任务。卫星提供广域态势和密钥来源，UAV提供可移动空中接入，岸基站承担计算和密钥管理，海上/水下平台提供数据采集、磁探测与可重构传播。",
            "这种分工会直接决定优化变量的来源：UAV位置形成轨迹变量，RIS控制形成相位矩阵，岸基调度形成二进制服务变量，QKD供给形成密钥预算，磁异常输出形成任务权重。系统图中的每条信息流都要能在模型中找到对应的变量或约束。",
            "本章因此先讲清楚节点、数据流和安全流的关系，再进入数学建模。这样做可以使后续公式不成为孤立符号，每一个符号都能回到具体的工程动作。",
        ],
        "本章作用：建立“系统节点—信息流—优化变量”的一一对应关系。",
    )
    page += 1

    s = new_slide()
    add_title(s, "SYSTEM LOGIC", "系统闭环：感知、调度、通信与安全回传")
    steps = [
        ("1 感知", "AUV/浮标发现海缆附近磁异常，产生风险优先级。"),
        ("2 建链", "UAV临时接入远海节点，RIS调整反射相位改善边缘链路。"),
        ("3 保护", "卫星QKD和岸基KMS为关键告警、控制指令和会话刷新分配密钥。"),
        ("4 优化", "统一求解轨迹、相位、功率、带宽、调度和密钥分配。"),
    ]
    for i, (t, b) in enumerate(steps):
        x = 0.82 + (i % 2) * 6.05
        y = 1.55 + (i // 2) * 1.82
        add_panel(s, x, y, 5.35, 1.25, RGBColor(252, 252, 249))
        add_text(s, t, x + 0.25, y + 0.16, 1.6, 0.32, 16, TEAL if i % 2 == 0 else CORAL, True)
        add_text(s, b, x + 0.25, y + 0.58, 4.75, 0.42, 13.5, INK)
    add_text(s, "这四步对应后续优化模型中的权重、信道、密钥约束和资源变量。", 0.95, 5.72, 10.8, 0.38, 17, MUTED)
    add_footer(s, page)
    page += 1

    image_slide(prs, page, "ARCHITECTURE", "空天地海一体化系统总体架构", ASSETS / "system_architecture_gpt.png", "该架构图给出四层节点之间的信息流与控制流。卫星层提供广域态势与QKD密钥，空中UAV层提供临时中继，地面岸基层完成优化求解与密钥管理，海上/水下层提供风电、海缆和磁探测数据。后续所有模型变量都从这些节点功能中抽象出来。", "System architecture diagram")
    notes.append("系统架构大图。")
    page += 1

    # 8 node table
    s = new_slide()
    add_title(s, "NODES", "系统节点功能与数学建模映射")
    rows = [
        ["节点", "功能", "进入模型的位置"],
        ["LEO卫星", "SAR遥感、QKD密钥下发", "密钥预算、灾后态势"],
        ["UAV", "移动中继、临时空中基站", "飞行轨迹、飞行能耗、速度约束"],
        ["岸基6G/KMS", "优化求解、密钥管理", "密钥分配、安全策略"],
        ["RIS平台", "被动波束赋形", "相位矩阵、2-bit硬件码本"],
        ["AUV/浮标", "磁异常检测与数据上浮", "任务权重、关键数据量"],
        ["电网中心", "告警处置与调度", "任务优先级与业务目标"],
    ]
    add_table(s, rows, 0.80, 1.45, 11.8, 4.65, col_widths=[0.20, 0.38, 0.42], font_size=11.2)
    add_footer(s, page)
    notes.append("节点功能表。")
    page += 1

    image_slide(prs, page, "DUAL FLOW", "业务数据流与安全控制流", ASSETS / "dual_flow_gpt.png", "图中区分了两类信息流：业务数据流主要经过UAV和RIS增强链路回传，安全控制流由QKD密钥预算保护。两者在调度层耦合，因为关键业务数据必须同时满足无线速率约束和密钥供给约束。", "Data and security-control flows")
    page += 1
    image_slide(prs, page, "QKD", "量子密钥分发与密钥预算约束", ASSETS / "qkd_mechanism_gpt.png", "该图解释QKD在系统中的工程角色：卫星链路主要周期性提供密钥材料，业务数据仍由UAV/RIS增强链路承担。密钥进入岸基密钥管理系统后，被分配给不同节点的关键数据流，因此在优化模型中表现为每个时隙的总密钥供给约束。", "QKD key management mechanism")
    page += 1
    image_slide(prs, page, "MAGNETIC SENSING", "量子磁力计海缆异常检测机制", ASSETS / "magnetometer_mechanism_gpt.png", "量子磁力计测量海缆附近的局部磁场扰动，并与WMM2025背景磁场进行差分。差分结果经过归一化后形成异常分数；该分数作为风险输入进入通信调度权重，使高风险海缆节点在资源受限时获得更高优先级。", "Quantum magnetic sensing mechanism")
    page += 1
    image_slide(prs, page, "RIS", "可重构智能反射超表面传播机制", ASSETS / "ris_mechanism_gpt.png", "RIS由大量可调反射单元组成。每个单元只能选择有限相位状态，本项目采用2-bit码本。通过相位对齐，RIS使直达链路和反射链路在UAV接收端尽量相干叠加，从而提高远海弱覆盖区域的等效信道增益。", "RIS propagation mechanism")
    page += 1

    transition_slide(
        prs,
        page,
        "三",
        "优化建模",
        "把系统图中的工程对象转成变量、目标函数和约束，形成可求解的能效-安全联合优化问题。",
        ["信道模型刻画RIS相干叠加。", "QKD模型刻画密钥供给瓶颈。", "磁异常模型把感知结果转成任务权重。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "MODEL INTRO",
        "优化建模导言：从工程限制到数学约束",
        [
            "建模的第一步是区分可控变量、外部输入和评价目标。可控变量包括轨迹、RIS相位、功率、带宽、调度和密钥分配；外部输入包括节点坐标、台风场景、磁异常背景和QKD供给；评价目标则是单位能耗下获得多少风险加权安全数据。",
            "每类工程限制都被写成约束：UAV机动能力对应速度与端点约束，RIS硬件对应单位模和离散码本约束，无线链路对应可达速率约束，QKD对应安全数据与密钥消耗约束，磁异常作为任务权重输入模型。",
            "由此得到的问题同时具有分式目标、二进制调度、离散相位、非凸轨迹和耦合资源分配。后续公式页将按“信道—密钥—优先级—能耗—完整目标”的顺序展开。",
        ],
        "本章作用：把工程系统转化为一个混合整数、非凸、分式优化问题。",
    )
    page += 1

    # 13 variables
    s = new_slide()
    add_title(s, "SYSTEM MODEL", "优化变量与约束体系")
    rows = [
        ["变量组", "物理含义", "主要约束"],
        ["UAV轨迹", "每个时隙的空中中继位置", "速度、初末点、电量"],
        ["RIS相位", "每个反射单元的相位状态", "单位模、2-bit离散码本"],
        ["功率与带宽", "各节点的无线资源分配", "功率上限、总带宽"],
        ["节点调度", "每时隙被服务的海上/水下节点", "二进制选择、容量限制"],
        ["安全数据量", "需要密钥保护的关键告警与控制数据", "速率和密钥共同限制"],
        ["密钥分配", "QKD密钥预算在节点之间的分配", "总密钥供给上限"],
    ]
    add_table(s, rows, 0.85, 1.45, 11.55, 4.70, col_widths=[0.22, 0.43, 0.35], font_size=12)
    add_text(s, "建模逻辑：系统架构中的每个工程动作都被映射为一个变量组；系统物理限制被写成约束；最终目标函数在这些约束下最大化加权安全数据能效。", 0.95, 6.24, 11.0, 0.42, 14.5, MUTED)
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "SYMBOLS I", "索引、集合与物理常量定义")
    symbol_rows = [
        (
            "离散索引",
            [r"k\in\mathcal{K}=\{1,\ldots,K\}", r"n\in\mathcal{N}=\{1,\ldots,N\}", r"m\in\mathcal{M}=\{1,\ldots,M\}"],
            "三类索引分别对应监测节点、离散时隙和RIS反射单元。后续所有求和、调度和相位更新都围绕这三个维度展开。",
            TEAL,
            "symbol_index_sets",
        ),
        (
            "相位码本",
            [r"\mathcal{Q}_{2\mathrm{bit}}=\{0,\pi/2,\pi,3\pi/2\}"],
            "2-bit码本给出RIS硬件允许选择的四个相位状态。连续相位先由对齐规则得到，再投影到该离散集合。",
            CORAL,
            "symbol_codebook",
        ),
        (
            "固定参数",
            [r"\tau,H,N_0>0,\quad P_{\max},B_{\max},V_{\max}>0", r"K^{\mathrm{QKD}}[n]\geq0,\quad \rho>0"],
            "这些参数定义系统边界：时隙长度、飞行高度、噪声水平、功率/带宽/速度上限，以及每个时隙可用的QKD密钥预算。",
            BLUE,
            "symbol_physical_parameters",
        ),
    ]
    for i, (heading, formulas, body, accent, fid) in enumerate(symbol_rows):
        y = 1.42 + i * 1.78
        add_text(s, heading, 0.92, y - 0.24, 2.30, 0.24, 12.5, accent, True)
        add_math_card(s, formulas, 0.85, y, 4.42, 1.30, accent, fid, fontsize=26 if i < 2 else 23)
        add_explanation_panel(s, heading + "说明", body, 5.55, y, 6.75, 1.30, accent, body_size=11.0)
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "SYMBOLS II", "决策变量、状态变量与收益变量定义")
    variable_rows = [
        (
            "几何与相位变量",
            [r"\mathbf{q}[n]=[x[n],y[n],H]^T,\quad \mathbf{\Theta}[n]=\mathrm{diag}(e^{j\theta_1[n]},\ldots,e^{j\theta_M[n]})"],
            "轨迹变量决定UAV空中中继的位置；相位矩阵决定RIS反射链路的相干叠加方式。",
            BLUE,
            "symbol_geometry_phase",
        ),
        (
            "调度与无线资源变量",
            [r"a_k[n]\in\{0,1\},\quad p_k[n]\geq0,\quad b_k[n]\geq0"],
            "调度变量表示节点是否被服务；功率和带宽变量表示该节点在当前时隙获得的无线资源。",
            TEAL,
            "symbol_schedule_resource",
        ),
        (
            "安全与任务收益变量",
            [r"s_k[n]\geq0,\quad \kappa_k[n]\geq0,\quad w_k[n]\geq0"],
            "安全数据量进入目标函数；密钥分配决定安全约束可行性；任务权重由磁异常和信息新鲜度生成。",
            CORAL,
            "symbol_security_priority",
        ),
        (
            "联合变量集合",
            [r"\mathbf{x}=\{\mathbf{q},\mathbf{\Theta},\mathbf{a},\mathbf{p},\mathbf{b},\mathbf{s},\boldsymbol{\kappa}\}"],
            "联合变量集合把可控量统一写入完整优化问题，便于表达目标函数、Dinkelbach变换和交替优化更新。",
            GOLD,
            "symbol_joint_variable",
        ),
    ]
    for i, (heading, formulas, body, accent, fid) in enumerate(variable_rows):
        y = 1.36 + i * 1.42
        add_text(s, heading, 0.92, y - 0.22, 2.75, 0.22, 12.0, accent, True)
        add_math_card(s, formulas, 0.85, y, 5.35, 1.02, accent, fid, fontsize=22 if i == 0 else 25)
        add_explanation_panel(s, heading + "说明", body, 6.50, y, 5.85, 1.02, accent, body_size=10.4)
    add_footer(s, page)
    page += 1

    # 14 channel
    s = new_slide()
    add_title(s, "CHANNEL", "RIS辅助等效信道与可达速率模型")
    add_formula(
        s,
        "等效信道与速率",
        [
            r"\mathbf{q}[n]=[x[n],y[n],H]^{T}",
            r"\mathbf{\Theta}[n]=\mathrm{diag}(e^{j\theta_1[n]},\ldots,e^{j\theta_M[n]})",
            r"h_k[n]=h_{kU}[n]+\mathbf{h}_{RU}^{H}[n]\mathbf{\Theta}[n]\mathbf{h}_{kR}",
            r"R_k[n]=b_k[n]\log_2\!\left(1+\frac{p_k[n]|h_k[n]|^2}{N_0 b_k[n]}\right)",
        ],
        0.80,
        1.55,
        6.15,
        4.80,
        TEAL,
        "channel_rate",
    )
    add_explanation_panel(
        s,
        "建模逻辑",
        "第一式定义UAV在离散时隙内的三维位置，其中高度在本项目中固定，水平坐标由轨迹优化更新。第二式定义RIS对角相位矩阵，每个对角元素对应一个反射单元的复反射系数，单位模约束由硬件无源反射特性决定。",
        7.15,
        1.58,
        5.00,
        2.20,
        TEAL,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "速率推导",
        "第三式把直达链路和RIS级联链路相加形成等效信道。级联项先经过节点到RIS链路，再乘以RIS相位矩阵，最后经过RIS到UAV链路。第四式采用Shannon容量形式：带宽决定频谱资源，功率和等效信道增益决定信噪比，噪声谱密度与带宽共同决定噪声功率。",
        7.15,
        4.05,
        5.00,
        2.05,
        CORAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "CHANNEL INTERPRETATION", "链路增益图与信道公式的关系")
    add_picture(s, FIG / "link_gain.png", 0.72, 1.35, 6.25, 4.95)
    add_explanation_panel(
        s,
        "图表内容",
        "左图比较不同传播控制策略下的远海链路速率，横轴是UAV与节点间的等效距离，纵轴是可达速率。无RIS曲线给出基线，随机RIS反映未进行相位控制时的反射效果，连续相位RIS给出理想上界，2-bit RIS给出硬件量化后的可实现结果。",
        7.25,
        1.35,
        4.95,
        2.20,
        TEAL,
        body_size=11.0,
    )
    add_explanation_panel(
        s,
        "结果含义",
        "图中的曲线随距离增加而下降，体现了信道公式中的路径损耗。相位优化后曲线上移，说明级联链路在UAV处相干叠加，提高了等效信道幅度。2-bit曲线接近连续相位上界，说明离散码本保留了主要波束赋形收益，是后续采用2-bit硬件约束的依据。",
        7.25,
        3.92,
        4.95,
        2.05,
        CORAL,
        body_size=11.0,
    )
    add_footer(s, page, text="Channel model validation")
    page += 1

    # 15 qkd constraints
    s = new_slide()
    add_title(s, "QKD CONSTRAINTS", "量子密钥供给与安全数据约束")
    add_formula(
        s,
        "QKD密钥约束",
        [
            r"0\leq s_k[n]\leq \tau R_k[n]",
            r"\rho s_k[n]\leq \tau\kappa_k[n]",
            r"\sum_k\kappa_k[n]\leq K^{\mathrm{QKD}}[n]",
            r"\rho=1\ \mathrm{for\ one\ time\ pad},\quad \rho<1\ \mathrm{for\ session\ refresh}",
        ],
        0.80,
        1.55,
        5.25,
        4.70,
        CORAL,
        "qkd_constraints",
    )
    add_explanation_panel(
        s,
        "约束含义",
        "第一式是通信容量上界：在一个时隙内被计入的安全数据不能超过无线链路可以承载的数据量。第二式是安全保护约束：每单位安全数据必须消耗相应密钥，密钥消耗系数越大，安全级别越接近一次一密。",
        6.65,
        1.58,
        5.35,
        2.10,
        CORAL,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "工程解释",
        "第三式给出同一时隙内所有节点共享的QKD密钥供给上限。它使优化问题不能只追求更高无线速率，还必须判断哪些数据值得消耗有限密钥。最后一式说明密钥消耗系数的安全含义：一次一密要求密钥量与数据量同阶，会话刷新则消耗更少但安全强度较低。",
        6.65,
        4.00,
        5.35,
        2.15,
        BLUE,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "QKD INTERPRETATION", "密钥供给曲线与安全数据瓶颈")
    add_picture(s, FIG / "qkd_secure_data.png", 0.72, 1.42, 6.30, 4.70)
    add_explanation_panel(
        s,
        "图表内容",
        "左图改变每个时隙可用的QKD密钥供给，记录最终能够安全回传的数据量。横轴表示密钥供给规模，纵轴表示安全数据吞吐；曲线前段上升较快，说明密钥是主要瓶颈。",
        7.25,
        1.40,
        4.95,
        2.05,
        CORAL,
        body_size=11.0,
    )
    add_explanation_panel(
        s,
        "与公式的关系",
        "当密钥预算较小时，第二式和第三式共同限制安全数据量；随着密钥供给增加，安全数据逐渐接近无线链路容量上界，此时第一式成为主要限制。该图说明系统瓶颈会在安全资源和通信资源之间转移，因此需要把QKD和无线资源放在同一个优化模型中。",
        7.25,
        3.80,
        4.95,
        2.30,
        TEAL,
        body_size=11.0,
    )
    add_footer(s, page, text="QKD key budget sensitivity")
    page += 1

    # 16 magnetic priority
    s = new_slide()
    add_title(s, "PRIORITY", "磁异常观测到任务优先级的数学映射")
    add_formula(
        s,
        "任务权重",
        [
            r"\Delta B_k[n]=B_{\mathrm{obs},k}[n]-B_{\mathrm{WMM2025},k}[n]",
            r"A_k[n]=\sigma\!\left(\frac{|\Delta B_k[n]|-\mu_B}{\sigma_B}\right)",
            r"w_k[n]=w_0+\lambda_A A_k[n]+\lambda_{\mathrm{AoI}}\mathrm{AoI}_k[n]",
            r"F(\mathbf{x})=\sum_n\sum_k w_k[n]s_k[n]",
        ],
        0.80,
        1.60,
        5.60,
        4.30,
        BLUE,
        "magnetic_priority",
    )
    add_explanation_panel(
        s,
        "从观测到异常分数",
        "第一式把量子磁力计观测值与WMM2025背景场相减，得到局部磁扰动。这样做可以去除大尺度地磁背景，使海缆附近的局部异常成为调度输入。第二式使用Sigmoid映射，把扰动幅值转为有界异常分数，避免少数极端观测值支配全部资源分配。",
        6.72,
        1.55,
        5.30,
        2.25,
        BLUE,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "从异常分数到收益函数",
        "第三式把异常分数和信息年龄共同写入任务权重，表示越异常、越久未更新的节点越值得优先服务。第四式把每个节点的安全数据量乘以任务权重后求和，形成分式目标中的收益部分。因此磁探测通过权重影响调度、资源和轨迹更新。",
        6.72,
        4.10,
        5.30,
        2.05,
        TEAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "PRIORITY INTERPRETATION", "磁异常优先级公式与热力图解读")
    add_picture(s, FIG / "magnetic_heatmap.png", 0.78, 1.35, 5.45, 5.15)
    add_explanation_panel(
        s,
        "公式链条",
        "该图对应公式页中的前两步：先从观测磁场中扣除WMM2025背景，再把局部扰动映射为异常分数。热力图先采样到各个监测节点，再转化为节点优先级并进入优化模型。",
        6.65,
        1.38,
        5.45,
        1.80,
        BLUE,
        body_size=10.0,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "左图颜色表示局部磁扰动强弱，黑色线段表示海缆路径，黄色星标表示设定的异常中心，散点为可服务节点。颜色越接近高扰动区域，节点风险权重越高。后续调度热力图和轨迹图中的资源集中现象，正是由这一风险权重进入收益函数后产生的。",
        6.65,
        3.50,
        5.45,
        2.25,
        TEAL,
        body_size=10.0,
    )
    add_footer(s, page, text="Magnetic anomaly to priority mapping")
    page += 1

    # 17 energy
    s = new_slide()
    add_title(s, "ENERGY MODEL", "系统总能耗模型")
    add_formula(
        s,
        "能耗模型",
        [
            r"E_{\mathrm{tot}}=E_{\mathrm{fly}}+E_{\mathrm{tx}}+E_{\mathrm{cir}}+E_{\mathrm{RIS}}",
            r"E_{\mathrm{fly}}=\sum_n\tau\left(P_0+\alpha\|\mathbf{v}[n]\|_2^2\right)",
            r"E_{\mathrm{tx}}=\sum_n\sum_k\tau p_k[n]",
            r"E_{\mathrm{RIS}}=N\tau M P_{\mathrm{elem}}",
        ],
        0.85,
        1.55,
        5.80,
        4.80,
        GOLD,
        "energy_model",
    )
    add_explanation_panel(
        s,
        "能耗分解",
        "第一式把系统总能耗分解为飞行、发射、电路和RIS硬件四部分。这样分解的意义在于区分可由轨迹改变的能耗、可由功率控制改变的能耗，以及由硬件规模近似固定决定的能耗。",
        7.05,
        1.55,
        5.05,
        2.00,
        GOLD,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "与优化目标的关系",
        "第二式采用速度平方项近似旋翼UAV机动代价，说明过度绕飞会快速增加能耗。第三式把每个时隙的发射功率积分为通信能耗。第四式说明RIS规模越大，反射孔径收益越强，但硬件单元也会带来额外能耗。因此最终目标从单纯最大化吞吐，提升为最大化单位能耗安全收益。",
        7.05,
        3.88,
        5.05,
        2.35,
        CORAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "ENERGY INTERPRETATION", "速度-能耗曲线与能耗模型的关系")
    add_picture(s, FIG / "uav_speed_energy.png", 0.72, 1.42, 6.35, 4.65)
    add_explanation_panel(
        s,
        "图表内容",
        "左图展示UAV速度变化引起的飞行能耗变化，并叠加优化轨迹中的速度序列。横轴表示时隙或速度采样点，纵轴表示相应能耗。曲线的上凸趋势来自公式页中的速度平方项。",
        7.25,
        1.42,
        4.95,
        2.00,
        GOLD,
        body_size=11.0,
    )
    add_explanation_panel(
        s,
        "结果含义",
        "该图解释为什么轨迹优化不能无限靠近所有高风险节点：若路径弯折过大，飞行能耗会抵消安全数据收益。SCA轨迹更新需要在异常节点收益、RIS覆盖收益和机动代价之间折中，这正是后续轨迹结果呈现“向高风险区域偏移但不过度绕飞”的原因。",
        7.25,
        3.78,
        4.95,
        2.25,
        TEAL,
        body_size=11.0,
    )
    add_footer(s, page, text="Energy model validation")
    page += 1

    # 18 complete optimization
    s = new_slide()
    add_title(s, "OBJECTIVE", "安全数据能效最大化问题")
    add_formula(
        s,
        "主优化问题",
        [
            r"\max_{\mathbf{x}\in\mathcal{X}}\ \eta(\mathbf{x})=\frac{F(\mathbf{x})}{G(\mathbf{x})}",
            r"F(\mathbf{x})=\sum_n\sum_k w_k[n]s_k[n]",
            r"G(\mathbf{x})=E_{\mathrm{fly}}+E_{\mathrm{tx}}+E_{\mathrm{cir}}+E_{\mathrm{RIS}}",
            r"\mathbf{x}=\{\mathbf{q},\mathbf{\Theta},\mathbf{p},\mathbf{b},\mathbf{a},\mathbf{s},\mathbf{\kappa}\}",
            r"\mathcal{X}:\ R,\ K^{\mathrm{QKD}},\ P,\ B,\ V_{\max},\ E_{\max},\theta_m\in\mathcal{Q}_{2\mathrm{bit}}",
        ],
        0.85,
        1.35,
        11.60,
        4.55,
        TEAL,
        "main_optimization",
    )
    add_text(s, "其中分子表示风险加权后的安全数据收益，分母表示系统总能耗，决策变量汇总轨迹、RIS相位、功率、带宽、调度、安全数据量和密钥分配。该问题同时包含分式目标、离散相位、二进制调度和非凸信道耦合。", 0.95, 6.12, 11.4, 0.62, 12.6, CORAL, True, leading=1.10)
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "OBJECTIVE DERIVATION", "目标函数构造逻辑")
    add_formula(
        s,
        "收益项",
        [
            r"F(\mathbf{x})=\sum_{n\in\mathcal{N}}\sum_{k\in\mathcal{K}}w_k[n]s_k[n]",
            r"w_k[n]=w_0+\lambda_A A_k[n]+\lambda_{\mathrm{AoI}}\mathrm{AoI}_k[n]",
        ],
        0.85,
        1.45,
        5.65,
        2.25,
        TEAL,
        "objective_benefit",
    )
    add_formula(
        s,
        "能耗与比值",
        [
            r"G(\mathbf{x})=E_{\mathrm{fly}}+E_{\mathrm{tx}}+E_{\mathrm{cir}}+E_{\mathrm{RIS}}",
            r"\eta(\mathbf{x})=F(\mathbf{x})/G(\mathbf{x})",
        ],
        6.80,
        1.45,
        5.65,
        2.25,
        GOLD,
        "objective_ratio",
    )
    add_explanation_panel(
        s,
        "推导说明",
        "收益项定义为风险加权后的安全数据量：磁异常越强、信息越过期，任务权重越高。分母采用系统总能耗，覆盖飞行、发射、电路和RIS硬件能耗。",
        0.85,
        4.05,
        5.65,
        1.90,
        TEAL,
    )
    add_explanation_panel(
        s,
        "优化含义",
        "最大化安全能效要求算法在额外能耗、密钥消耗和风险收益之间做折中。该比值目标也是后续引入Dinkelbach分式规划的直接原因。",
        6.80,
        4.05,
        5.65,
        1.90,
        GOLD,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "NONCONVEXITY", "非凸性来源与算法分解依据")
    hard_points = [
        ("分式目标", "安全数据收益除以总能耗，需要 Dinkelbach 参数化。"),
        ("离散变量", "节点调度和RIS 2-bit相位引入组合搜索。"),
        ("非凸信道", "UAV位置进入距离、路径损耗和RIS反射链路。"),
        ("资源耦合", "功率、带宽、密钥和安全数据量互相限制。"),
        ("轨迹约束", "速度、电量和初末点限制让每个时隙不独立。"),
    ]
    for i, (t, b) in enumerate(hard_points):
        y = 1.42 + i * 0.80
        add_text(s, f"{i + 1}", 0.95, y, 0.35, 0.30, 15, CORAL, True, font=FONT_EN)
        add_text(s, t, 1.45, y - 0.02, 1.65, 0.32, 15.5, INK, True)
        add_text(s, b, 3.25, y - 0.02, 8.25, 0.34, 14.5, MUTED)
    add_text(s, "因此算法不能只做单一模块优化，而要把分式规划、交替优化、凸子问题和轨迹近似组合起来。", 0.95, 6.10, 10.8, 0.38, 17, INK, True)
    add_footer(s, page)
    page += 1

    transition_slide(
        prs,
        page,
        "四",
        "Q-RIS-UAV-AO算法",
        "算法的目标是把不可直接求解的混合整数非凸分式问题，拆成一组有物理含义、可复现的子问题。",
        ["Dinkelbach 处理分式安全能效。", "AO 分别更新调度、资源、RIS相位和UAV轨迹。", "CVXPY求解资源分配凸子问题，SCA更新轨迹。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "ALGORITHM INTRO",
        "算法导言：Dinkelbach-AO-SCA分解思路",
        [
            "完整问题不能直接交给一个通用凸优化器求解，原因在于分式目标、离散调度、RIS离散相位和轨迹非凸项同时存在。若直接枚举，组合复杂度会随着节点数、时隙数和RIS单元数快速增长。",
            "本项目采用分层分解：外层Dinkelbach把安全数据能效比值转化为参数化差分问题；内层交替优化固定其余变量，依次更新调度、资源、RIS相位和UAV轨迹；连续资源子问题由CVXPY求解，轨迹子问题用逐次凸近似构造局部可解更新。",
            "该分解强调可解释、可复现和约束可行。每一步都有明确的物理意义，并且都能在实验图中看到对应效果。",
        ],
        "本章作用：说明为什么联合优化问题要被拆成四类可计算子问题。",
    )
    page += 1

    image_slide(prs, page, "DECOMPOSITION", "混合整数非凸分式问题的分解框架", ASSETS / "optimization_theory_gpt.png", "该图展示算法分解的理论逻辑：外层使用Dinkelbach方法处理安全能效分式目标，内层使用交替优化（Alternating Optimization, AO）分别更新调度、资源、RIS相位和UAV轨迹。每个子问题都对应前一章的变量组，因此算法步骤与数学模型是一一对应的。", "Optimization decomposition diagram")
    page += 1

    # 20 Dinkelbach derivation
    s = new_slide()
    add_title(s, "DINKELBACH", "Dinkelbach分式规划变换")
    add_formula(
        s,
        "推导",
        [
            r"\Phi(\eta)=\max_{\mathbf{x}\in\mathcal{X}}\{F(\mathbf{x})-\eta G(\mathbf{x})\}",
            r"\mathbf{x}^{(t)}=\arg\max_{\mathbf{x}\in\mathcal{X}}\{F(\mathbf{x})-\eta^{(t)}G(\mathbf{x})\}",
            r"\eta^{(t+1)}=\frac{F(\mathbf{x}^{(t)})}{G(\mathbf{x}^{(t)})}",
            r"|\Phi(\eta^{(t)})|\leq\epsilon",
        ],
        0.85,
        1.45,
        5.85,
        4.85,
        CORAL,
        "dinkelbach",
    )
    add_explanation_panel(
        s,
        "变换目标",
        "原始目标是收益与能耗之比。Dinkelbach方法引入一个能效参数，把比值最大化改写为参数化差分最大化。若当前参数低于最优能效，则存在可行解使差分值为正；若参数高于最优能效，差分值不能为正。",
        7.05,
        1.52,
        5.05,
        2.18,
        CORAL,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "迭代规则",
        "第二式在当前参数下求解差分子问题，得到本轮联合变量。第三式用本轮收益与能耗之比更新参数，相当于把下一次搜索中心移动到当前可行解的安全能效。第四式给出停止准则：当差分残差足够接近零时，当前参数即为分式目标的近似最优值。",
        7.05,
        4.05,
        5.05,
        2.08,
        TEAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "DINKELBACH INTERPRETATION", "Dinkelbach迭代曲线与分式目标的关系")
    add_picture(s, FIG / "dinkelbach_eta_gap.png", 0.72, 1.38, 5.70, 4.25)
    add_explanation_panel(
        s,
        "数学原理",
        "原问题为 max 𝑭(𝒙)/𝑮(𝒙)。当 η 低于最优比值时，存在某个可行解使 𝑭(𝒙)-η𝑮(𝒙)>0；当 η 高于最优比值时，所有可行解的差分值均不会为正。因此最优比值 η⋆ 等价于方程 Φ(η⋆)=0 的根。",
        6.75,
        1.35,
        5.45,
        1.80,
        CORAL,
        body_size=9.8,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "曲线记录每一轮外层迭代的 η 与差分残差。η 上升表示当前联合变量产生了更高的单位能耗安全收益；残差下降表示差分问题逐步接近零点。该图与公式页的停止准则 |Φ(η)|≤ε 直接对应，用于验证分式规划外层迭代的数值稳定性。",
        6.75,
        3.45,
        5.45,
        2.10,
        TEAL,
        body_size=9.8,
    )
    add_footer(s, page, text="Dinkelbach residual and secure energy efficiency")
    page += 1

    # 21 schedule
    s = new_slide()
    add_title(s, "SCHEDULING", "节点调度子问题与优先指数")
    add_formula(
        s,
        "调度优先指数",
        [
            r"\mathrm{Score}_k[n]=w_k[n]\hat{R}_k[n]+\beta D_{\mathrm{rem},k}[n]-\mu\mathrm{AoI}_k[n]",
            r"a_k[n]\in\{0,1\}",
            r"\sum_k a_k[n]\leq A_{\max}",
            r"a_k[n]=1\Rightarrow k\in\mathcal{N}_{\mathrm{cover}}(\mathbf{q}[n])",
        ],
        0.85,
        1.55,
        5.95,
        4.60,
        TEAL,
        "scheduling_score",
    )
    add_explanation_panel(
        s,
        "优先指数构成",
        "第一式定义调度优先指数，由三类因素组成：风险加权链路收益、剩余任务需求和信息年龄惩罚。风险权重使磁异常节点更容易被服务，速率估计保证被选节点具备可用链路，剩余任务项防止节点长期得不到传输机会。",
        7.05,
        1.55,
        5.00,
        2.28,
        TEAL,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "可行性约束",
        "第二式规定调度变量只能取服务或不服务两种状态。第三式限制一个时隙中最多同时服务的节点数量，反映UAV接入和解调能力有限。第四式要求被服务节点必须位于当前UAV位置对应的覆盖集合内，使调度决策与轨迹几何位置保持一致。",
        7.05,
        4.12,
        5.00,
        2.02,
        BLUE,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "SCHEDULING INTERPRETATION", "节点调度热力图与优先指数的关系")
    add_picture(s, FIG / "secure_data_heatmap.png", 0.72, 1.35, 5.85, 4.55)
    add_explanation_panel(
        s,
        "运算规则",
        "在每个时隙 𝑛，算法先计算所有节点的 𝑺coreₖ[𝑛]。其中 𝑤ₖ[𝑛] 提高异常节点收益，R̂ₖ[𝑛] 反映当前链路质量，𝐷rem,ₖ[𝑛] 避免任务长期未完成，AoIₖ[𝑛] 控制信息陈旧度。随后在覆盖集合和并发数量约束下选择 𝑎ₖ[𝑛]=1 的节点。",
        6.85,
        1.35,
        5.30,
        2.05,
        TEAL,
        body_size=9.8,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "左图横轴为时隙、纵轴为节点，颜色表示最终分配的安全数据量 𝑠ₖ[𝑛]。热力图呈现稀疏条带，说明调度会在每个时隙把资源集中给满足覆盖、密钥和链路条件的高优先级节点。",
        6.85,
        3.72,
        5.30,
        1.85,
        BLUE,
        body_size=9.8,
    )
    add_footer(s, page, text="Scheduling index and secure-data heatmap")
    page += 1

    # 22 cvxpy
    s = new_slide()
    add_title(s, "CVXPY SUBPROBLEM", "功率、密钥与安全数据量的凸优化子问题")
    add_formula(
        s,
        "CVXPY资源分配",
        [
            r"\max_{\mathbf{p},\mathbf{s},\mathbf{\kappa}}\ \sum_k w_k s_k-\eta\!\left(E_{\mathrm{fix}}+c\sum_k p_k\right)",
            r"s_k\leq b_k\log_2(1+\gamma_k p_k)",
            r"\rho s_k\leq\kappa_k,\quad \sum_k\kappa_k\leq K^{\mathrm{QKD}}",
            r"0\leq p_k\leq P_{\max}a_k",
        ],
        0.80,
        1.45,
        5.85,
        4.85,
        BLUE,
        "cvxpy_resource",
    )
    add_explanation_panel(
        s,
        "子问题目标",
        "当UAV轨迹、RIS相位、带宽和调度固定后，资源子问题只需要决定发射功率、安全数据量和密钥分配。第一式中的收益项奖励高权重安全数据，惩罚项来自Dinkelbach差分目标，用于抑制为了少量吞吐而消耗过多功率。",
        6.95,
        1.52,
        5.20,
        2.20,
        BLUE,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "凸性与约束",
        "第二式是固定信道增益后的速率上界，关于发射功率为凹函数；第三式是密钥消耗和总密钥预算约束；第四式把功率上界与调度开关绑定，使未被调度节点不能发射。该子问题满足标准凸优化求解器可处理的形式，因此用CVXPY稳定求解。",
        6.95,
        4.05,
        5.20,
        2.12,
        TEAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "RESOURCE INTERPRETATION", "CVXPY资源分配结果与凸子问题解读")
    add_picture(s, FIG / "power_allocation_heatmap.png", 0.70, 1.35, 5.75, 4.30)
    add_explanation_panel(
        s,
        "子问题含义",
        "AO固定其他变量后，资源子问题只决定 𝑝ₖ、𝑏ₖ、𝑠ₖ 与 κₖ。目标项 Σ𝑤ₖ𝑠ₖ 奖励高风险安全数据，惩罚项 η·𝑐Σ𝑝ₖ 抑制过度发射功率。由于 log₂(1+γₖ𝑝ₖ) 对 𝑝ₖ 凹，最大化凹收益并满足线性约束可由CVXPY求解。",
        6.75,
        1.35,
        5.45,
        2.05,
        BLUE,
        body_size=9.6,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "左图颜色表示不同时隙和节点上的发射功率。功率会在链路增益 γₖ、节点权重 𝑤ₖ、密钥预算 κₖ 与调度开关 𝑎ₖ 之间折中。该结果说明资源分配子问题承担了把“高优先级”转化为实际通信能量的作用。",
        6.75,
        3.72,
        5.45,
        1.95,
        TEAL,
        body_size=9.6,
    )
    add_footer(s, page, text="CVXPY resource allocation subproblem")
    page += 1

    # 23 RIS phase
    s = new_slide()
    add_title(s, "RIS PHASE", "RIS相位对齐与2-bit码本投影")
    add_formula(
        s,
        "相位投影",
        [
            r"\theta_m^\star=-\arg(h_{RU,m}^{\ast}h_{kR,m})+\arg(h_{kU})",
            r"\mathcal{Q}_{2\mathrm{bit}}=\{0,\pi/2,\pi,3\pi/2\}",
            r"\theta_m\leftarrow\Pi_{\mathcal{Q}_{2\mathrm{bit}}}(\theta_m^\star)",
            r"\mathbf{\Theta}=\mathrm{diag}(e^{j\theta_1},\ldots,e^{j\theta_M})",
        ],
        0.80,
        1.45,
        5.85,
        4.85,
        GOLD,
        "ris_phase_projection",
    )
    add_explanation_panel(
        s,
        "连续相位对齐",
        "第一式给出第m个RIS单元的理想连续相位。其目标是抵消节点到RIS、RIS到UAV两段传播中的相位差，并与直达链路参考相位对齐，使各反射单元贡献在UAV接收端尽量同相相加。",
        7.05,
        1.52,
        5.05,
        2.18,
        GOLD,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "离散投影",
        "第二式定义2-bit硬件码本，只有四个可选相位状态。第三式把连续最优相位投影到最近码字，体现实际RIS单元不能连续调相的约束。第四式将所有单元相位组装成对角矩阵，供等效信道公式调用。",
        7.05,
        4.02,
        5.05,
        2.10,
        TEAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "RIS PHASE INTERPRETATION", "RIS相位编码图与2-bit投影规则")
    add_picture(s, FIG / "ris_matlab_phase_coding.png", 0.78, 1.32, 5.35, 4.50)
    add_explanation_panel(
        s,
        "相位对齐原理",
        "级联链路由节点到RIS的相位、RIS单元反射相位和RIS到UAV的相位共同决定。连续最优 θₘ⋆ 的作用是抵消前后两段传播相位差，使各单元反射信号在UAV处尽量同相叠加，从而提升等效信道幅度。",
        6.65,
        1.35,
        5.45,
        1.90,
        GOLD,
        body_size=9.8,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "左图展示16×16单元的2-bit相位编码矩阵，颜色对应0、π/2、π和3π/2四种相位状态。相邻区域呈现条纹式变化，反映了指向目标角时需要沿阵列孔径逐渐补偿的传播相位梯度。该图把公式中的投影θₘ←Π𝒬(θₘ⋆)转化为可观察的硬件编码结果。",
        6.65,
        3.55,
        5.45,
        2.15,
        TEAL,
        body_size=9.6,
    )
    add_footer(s, page, text="RIS 2-bit phase projection")
    page += 1

    # 24 UAV SCA
    s = new_slide()
    add_title(s, "UAV SCA", "UAV轨迹的逐次凸近似更新")
    add_formula(
        s,
        "轨迹更新",
        [
            r"R_k(\mathbf{q}[n])\geq R_k(\mathbf{q}^{(r)}[n])+\nabla R_k(\mathbf{q}^{(r)}[n])^T(\mathbf{q}[n]-\mathbf{q}^{(r)}[n])",
            r"\|\mathbf{q}[n+1]-\mathbf{q}[n]\|_2\leq V_{\max}\tau",
            r"\mathbf{q}[1]=\mathbf{q}_0,\quad \mathbf{q}[N]=\mathbf{q}_F",
            r"\mathbf{q}^{(r+1)}=\mathbf{q}^{(r)}+\alpha\Delta\mathbf{q}",
        ],
        0.80,
        1.45,
        6.15,
        4.85,
        TEAL,
        "uav_sca",
    )
    add_explanation_panel(
        s,
        "一阶近似",
        "第一式是在第r轮当前轨迹处对速率函数构造的一阶下界。由于UAV位置会同时改变直达距离、RIS级联距离和路径损耗，原速率函数关于轨迹通常非凸；用局部下界替代后，可以得到可求解的近似更新方向。",
        7.15,
        1.52,
        4.95,
        2.18,
        TEAL,
        body_size=10.9,
    )
    add_explanation_panel(
        s,
        "可行更新",
        "第二式限制相邻时隙之间的最大位移，对应UAV速度上限；第三式固定起点和终点，保证任务路径可执行；第四式用步长控制从当前轨迹向近似问题解移动，避免一次更新过大导致线性化失效。该过程重复执行，形成逐次凸近似轨迹优化。",
        7.15,
        4.02,
        4.95,
        2.12,
        CORAL,
        body_size=10.9,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "TRAJECTORY INTERPRETATION", "UAV轨迹图与SCA更新规则")
    add_picture(s, FIG / "trajectory.png", 0.95, 1.20, 4.85, 5.60)
    add_explanation_panel(
        s,
        "更新目标",
        "在第𝑟轮迭代中，SCA将非凸速率 𝑅ₖ(𝒒[𝑛]) 用当前轨迹 𝒒⁽ʳ⁾[𝑛] 处的一阶下界替代。这样得到的近似问题会鼓励轨迹向高权重、高链路收益区域移动，但每一步仍满足 ‖𝒒[𝑛+1]-𝒒[𝑛]‖₂≤𝑉maxτ 与初末点约束。",
        6.45,
        1.35,
        5.65,
        2.05,
        TEAL,
        body_size=9.6,
    )
    add_explanation_panel(
        s,
        "图表解读",
        "左图灰色虚线为初始巡检路径，红色曲线为优化路径，颜色散点表示磁异常优先级。优化后轨迹向黄色异常星标和高权重节点附近弯曲，同时没有直接追逐所有节点，说明SCA更新在任务收益、飞行代价和RIS覆盖之间进行了局部折中。",
        6.45,
        3.72,
        5.65,
        1.95,
        CORAL,
        body_size=9.6,
    )
    add_footer(s, page, text="SCA trajectory update and optimized path")
    page += 1

    image_slide(prs, page, "ALGORITHM", "Q-RIS-UAV-AO算法流程", ASSETS / "algorithm_flow_gpt.png", "该流程图给出一次完整迭代：首先根据当前变量计算安全能效；随后在固定其余变量的条件下依次更新调度、功率与密钥、RIS相位和UAV轨迹；最后用Dinkelbach准则更新分式目标参数。算法停止条件由残差和最大迭代次数共同控制。", "Algorithm workflow")
    page += 1

    s = new_slide()
    add_title(s, "ALGORITHM OUTPUTS", "算法输出与系统可执行控制量")
    outputs = [
        ("UAV轨迹", "每个时隙的位置序列，用于灾后临时中继巡航。"),
        ("RIS相位码本", "每个时隙的2-bit相位矩阵，用于远海链路增强。"),
        ("功率与带宽", "在总资源约束下优先给高价值节点。"),
        ("密钥分配", "把有限QKD密钥分配给关键安全数据。"),
        ("安全数据指标", "输出安全能效、约束违反量和收敛曲线。"),
    ]
    for i, (t, b) in enumerate(outputs):
        x = 0.85 + (i % 2) * 5.95
        y = 1.45 + (i // 2) * 1.10
        add_panel(s, x, y, 5.25, 0.82, RGBColor(252, 252, 249))
        add_text(s, t, x + 0.22, y + 0.13, 1.55, 0.26, 14.5, TEAL, True)
        add_text(s, b, x + 1.78, y + 0.13, 3.10, 0.30, 12.5, INK)
    add_text(s, "复杂度控制策略：RIS相位使用投影量化避免指数级搜索，资源分配交给凸优化器保证数值稳定，轨迹只在当前点附近做逐次凸近似更新。这样既保留非凸联合优化的主要结构，又能在课程项目规模内完整复现。", 0.95, 5.55, 11.0, 0.70, 13.8, MUTED, leading=1.12)
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "PSEUDOCODE", "Algorithm 1：Q-RIS-UAV-AO伪代码")
    add_panel(s, 0.70, 1.35, 7.35, 5.55, RGBColor(22, 32, 42), RGBColor(22, 32, 42))
    pseudo_lines = [
        "Algorithm 1  Q-RIS-UAV-AO",
        "Input : node set 𝒦, slots 𝒩, RIS codebook 𝒬₂bit,",
        "        QKD supply Kᴽᴷᴰ[n], budgets Pₘₐₓ, Bₘₐₓ, Eₘₐₓ",
        "Output: 𝒒[n], 𝜣[n], aₖ[n], pₖ[n], bₖ[n], sₖ[n], κₖ[n]",
        "",
        "1: initialize 𝒒⁽⁰⁾, 𝜣⁽⁰⁾, 𝒂⁽⁰⁾, η⁽⁰⁾",
        "2: for t = 0,1,...,Tₘₐₓ do",
        "3:     compute channel hₖ[n] and priority wₖ[n]",
        "4:     update scheduling aₖ[n] by 𝑺coreₖ[n]",
        "5:     solve {pₖ,bₖ,sₖ,κₖ} by convex resource allocation",
        "6:     align RIS phase θₘ* and project it to 𝒬₂bit",
        "7:     update UAV trajectory 𝒒 by SCA lower-bound step",
        "8:     η⁽ᵗ⁺¹⁾ = 𝑭(𝒙⁽ᵗ⁾) / 𝑮(𝒙⁽ᵗ⁾)",
        "9:     stop if |𝑭(𝒙⁽ᵗ⁾)-η⁽ᵗ⁾𝑮(𝒙⁽ᵗ⁾)| ≤ ε",
        "10: end for",
    ]
    for i, line in enumerate(pseudo_lines):
        y = 1.58 + i * 0.31
        color = GOLD if i == 0 else RGBColor(235, 240, 244)
        size = 12.0 if i == 0 else 9.4
        add_text(s, line, 1.00, y, 6.85, 0.24, size, color, i == 0, font=FONT_MONO)
    add_explanation_panel(
        s,
        "读法说明",
        "该页给出论文伪码格式：输入集合、预算和码本，输出轨迹、相位、调度、资源和密钥。第4至7行对应四个交替优化子问题，第8至9行对应Dinkelbach外层参数更新与停止判据。",
        8.35,
        1.40,
        3.95,
        2.35,
        TEAL,
        body_size=10.0,
    )
    add_explanation_panel(
        s,
        "计算规模",
        "直接枚举𝒂和𝜣会导致组合复杂度随|𝒦|、|𝒩|和RIS单元数指数增长。伪代码通过“离散启发式 + 凸资源分配 + 局部轨迹近似”把求解转化为可复现实验规模。",
        8.35,
        4.05,
        3.95,
        2.25,
        CORAL,
        body_size=10.0,
    )
    add_footer(s, page)
    page += 1

    s = new_slide()
    add_title(s, "PSEUDOCODE", "算法伪代码与停止准则")
    add_panel(s, 0.85, 1.40, 6.05, 4.95, DEEP, DEEP)
    steps = [
        "初始化UAV轨迹、RIS相位、调度变量和能效参数。",
        "固定当前变量，计算节点权重、信道增益和可服务集合。",
        "更新调度，再用CVXPY求解功率、密钥和安全数据量。",
        "按相位对齐规则更新RIS相位，并投影到2-bit码本。",
        "用逐次凸近似更新UAV轨迹，随后更新Dinkelbach参数。",
        "若差分残差低于阈值或达到最大迭代次数，则停止。",
    ]
    for i, step_text in enumerate(steps):
        y = 1.75 + i * 0.62
        add_text(s, f"{i + 1}", 1.12, y, 0.32, 0.26, 11.5, GOLD, True, font=FONT_EN)
        add_text(s, step_text, 1.55, y - 0.02, 4.85, 0.34, 12.8, WHITE)
    add_formula(
        s,
        "停止准则",
        [
            r"\Phi(\eta^{(t)})=F(\mathbf{x}^{(t)})-\eta^{(t)}G(\mathbf{x}^{(t)})",
            r"|\Phi(\eta^{(t)})|\leq\epsilon\quad\mathrm{or}\quad t=T_{\max}",
        ],
        7.25,
        1.55,
        4.95,
        2.45,
        CORAL,
        "algorithm_stop",
    )
    add_explanation_panel(
        s,
        "数值实现说明",
        "伪代码中的每一步都对应一个已经定义的子问题。调度和相位更新提供可行的离散决策，CVXPY子问题保证连续资源分配的数值稳定，轨迹更新保持速度和初末点约束，停止准则监控分式规划残差。",
        7.25,
        4.35,
        4.95,
        1.70,
        CORAL,
    )
    add_footer(s, page)
    page += 1

    transition_slide(
        prs,
        page,
        "五",
        "仿真实验与分析",
        "实验按照“链路机制—资源分配—轨迹行为—安全约束—鲁棒性”的顺序展开，逐步验证每个模块的作用。",
        ["验证RIS相位控制对远海链路增益的影响。", "验证联合优化算法的收敛性及其相对基线优势。", "检验QKD密钥、磁异常优先级和扰动条件下的系统行为。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "EXPERIMENT INTRO",
        "实验导言：从物理机制到优化结果的证据链",
        [
            "实验部分按照模型中的因果链组织，以替代结果罗列式展示。首先验证RIS相位编码确实带来链路和远场主瓣收益；随后验证资源分配、QKD约束和UAV轨迹是否按照模型预期改变；最后用收敛、消融、鲁棒性和约束检查说明算法行为稳定且可行。",
            "每一类图表都对应一个数学模块：链路速率图对应信道公式，密钥曲线对应QKD约束，轨迹图对应SCA更新，热力图对应资源分配，Pareto曲线对应能耗与安全数据之间的权衡。",
            "因此，实验展示围绕一个连续问题展开：引入量子感知、量子密钥和RIS传播控制后，联合优化是否能把有限资源更集中地用于高风险安全数据。",
        ],
        "本章作用：把公式模型转化为可观察、可复现、可比较的仿真证据。",
    )
    page += 1

    image_slide(prs, page, "EXPERIMENT DESIGN", "实验设计与评价指标体系", ASSETS / "experiment_design_gpt.png", "该图给出实验逻辑链。首先用链路与RIS仿真验证物理机制，再用收敛曲线和消融实验验证算法有效性，随后用QKD、磁异常和鲁棒性实验检验系统约束下的行为。所有图表均由脚本生成，并使用固定随机种子保证可复现。", "Experiment design overview")
    page += 1

    # 27 parameter table
    s = new_slide()
    add_title(s, "PARAMETERS", "实验参数与复现设置")
    rows = [
        ["模块", "设置", "说明"],
        ["场景", "DeepOWT东海聚类24节点", "真实海上风电基础设施"],
        ["灾害", "NOAA IBTrACS台风轨迹", "真实灾害背景"],
        ["磁场", "WMM2025背景 + 合成偶极异常", "量子磁探测闭环"],
        ["通信", "5.8 GHz，256单元RIS，2-bit相位", "工程可实现参数"],
        ["优化", "Dinkelbach + AO + CVXPY + SCA", "课程优化理论主线"],
        ["复现", "Python 3.12, fixed seed 20260513", "一键运行run_all_v2.ps1"],
    ]
    add_table(s, rows, 0.85, 1.45, 11.60, 4.90, col_widths=[0.18, 0.43, 0.39], font_size=11.5)
    add_text(s, "参数设计遵循两条原则：空间场景尽量来自公开数据，通信与优化部分使用可控仿真以便分析算法机制。RIS规模、相位量化和载频选取保持在常见系统级论文可讨论范围内。", 0.95, 6.22, 11.0, 0.46, 13.5, MUTED)
    add_footer(s, page)
    page += 1

    image_slide(prs, page, "RESULT 1", "RIS相位优化对链路速率的提升", FIG / "link_gain.png", "该图比较无RIS、随机RIS、连续相位RIS和2-bit量化RIS在不同距离下的速率。速率随距离增大整体下降，符合路径损耗规律；相位优化后曲线明显上移，说明RIS能够补偿远海边缘覆盖损失。2-bit量化曲线接近连续相位上界，表明低比特硬件仍能获得主要波束赋形收益。")
    page += 1
    image_slide(prs, page, "RESULT 2", "RIS单元数敏感性分析", FIG / "ris_elements_rate.png", "该图考察RIS单元数量对平均速率的影响。随着单元数增加，等效反射孔径和可合成信号能量增大，因此平均速率单调提高；但高规模区域的增益斜率开始变缓，说明RIS规模存在边际收益。该结果为后续选择256单元作为折中规模提供依据。")
    page += 1
    image_slide(prs, page, "RESULT 3", "RIS相位量化精度对远场波束的影响", FIG / "ris_quantization_comparison.png", "该图比较1-bit、2-bit和3-bit相位量化下的归一化远场方向图。1-bit量化造成明显主瓣损失和旁瓣抬升；2-bit已经接近3-bit曲线，说明四状态相位编码足以捕获主要相位对齐收益。工程上，2-bit在控制复杂度和波束质量之间更均衡。")
    page += 1
    image_slide(prs, page, "RESULT 4", "MATLAB远场波束赋形仿真", FIG / "ris_matlab_farfield.png", "该图由MATLAB R2024a运行阵列因子脚本得到。仿真目标角设置为30度，连续相位和2-bit相位两条曲线的主瓣均指向目标方向附近，说明离散相位投影没有破坏主瓣指向能力。该页回应超表面/RIS仿真要求中的远场波束赋形部分。")
    page += 1

    image_slide(prs, page, "RESULT 5", "RIS单元四状态相位响应", FIG / "ris_matlab_unit_response.png", "该图模拟2-bit RIS单元在频率扫描下的相位响应。四条曲线对应0度、90度、180度和270度相位状态，显示不同编码状态在工作频段内保持可区分性。它为系统级2-bit相位码本提供单元响应层面的支撑。")
    page += 1

    s = new_slide()
    add_title(s, "RESULT 6", "二维远场方向图与目标角扫描")
    add_picture(s, FIG / "ris_2d_farfield_heatmap.png", 0.72, 1.45, 5.45, 4.15)
    add_picture(s, FIG / "ris_beam_scan_heatmap.png", 6.35, 1.45, 5.45, 4.15)
    add_text(s, "左图给出二维方位-俯仰远场热力图，主瓣集中在目标方向附近；右图改变目标扫描角，主瓣位置随相位编码同步移动。两张图共同说明RIS能够通过相位配置实现可重构波束控制。", 0.95, 5.95, 11.0, 0.65, 13.2, MUTED, leading=1.12)
    add_footer(s, page, text="RIS 2D far-field and beam scan")
    page += 1
    s = new_slide()
    add_title(s, "RESULT 7", "风险感知UAV轨迹优化结果")
    add_panel(s, 0.85, 1.58, 4.85, 4.70, RGBColor(252, 252, 249))
    add_text(s, "轨迹解释", 1.10, 1.88, 2.0, 0.34, 16, TEAL, True)
    add_bullets(
        s,
        [
            "图中灰色虚线为初始巡检轨迹，红色曲线为优化后轨迹。",
            "优化轨迹在高权重磁异常节点附近偏移，提高关键数据回传机会。",
            "轨迹仍满足速度与初末点约束，并兼顾RIS覆盖带来的链路收益。",
        ],
        1.12,
        2.35,
        4.25,
        2.85,
        size=12.8,
        accent=TEAL,
    )
    add_picture(s, FIG / "trajectory.png", 6.10, 1.35, 6.25, 5.45)
    add_footer(s, page, text="Optimized UAV trajectory and magnetic priority")
    page += 1
    image_slide(prs, page, "RESULT 8", "RIS-UAV联合覆盖热力图", FIG / "coverage_heatmap.png", "该热力图展示空间位置上的等效覆盖强度，颜色越亮表示在RIS辅助下的可达速率越高；橙色曲线为优化后的UAV轨迹，青色菱形为RIS平台位置。优化轨迹经过覆盖较强区域并靠近高风险节点，说明轨迹优化同时考虑通信收益和任务优先级。")
    page += 1
    image_slide(prs, page, "RESULT 9", "Dinkelbach外层迭代收敛性", FIG / "dinkelbach_eta_gap.png", "该图记录外层迭代中的安全能效参数和差分残差。能效参数逐步上升，残差下降后进入稳定区间，说明参数化差分问题正在逼近分式目标的零点。该结果验证了Dinkelbach外层更新在本仿真场景中的数值可行性。")
    page += 1
    image_slide(prs, page, "RESULT 10", "交替优化算法与消融基线收敛对比", FIG / "convergence.png", "该图比较Q-RIS-UAV-AO与只优化RIS、只优化UAV、随机策略等基线。联合优化曲线收敛后保持最高安全能效，说明轨迹、相位和资源分配之间存在互补关系；单独优化某一模块只能获得部分收益。")
    page += 1
    image_slide(prs, page, "RESULT 11", "安全数据分配热力图", FIG / "secure_data_heatmap.png", "该热力图由CVXPY资源分配子问题输出。横轴为时隙，纵轴为节点编号，颜色表示安全数据量。结果呈现稀疏的时隙-节点结构，说明在密钥预算和调度容量约束下，算法会把安全数据集中分配给高权重、链路条件较好的节点。")
    page += 1

    s = new_slide()
    add_title(s, "RESULT 12", "QKD密钥分配与安全数据瓶颈")
    add_picture(s, FIG / "key_allocation_heatmap.png", 0.72, 1.55, 5.85, 4.65)
    add_picture(s, FIG / "qkd_secure_data.png", 6.72, 1.55, 5.85, 4.65)
    add_text(s, "左图给出每个时隙的密钥分配结构，说明CVXPY会把有限QKD密钥优先分配给高价值节点；右图改变密钥供给规模，安全数据量随密钥增加而提升，并在链路吞吐上限附近趋于饱和。该结果说明系统瓶颈会在密钥约束和无线链路容量之间转移。", 0.95, 6.12, 11.0, 0.58, 13.0, MUTED, leading=1.12)
    add_footer(s, page, text="QKD key allocation and secure-data response")
    page += 1
    s = new_slide()
    add_title(s, "RESULT 14", "磁异常检测与风险优先级映射")
    add_panel(s, 0.85, 1.58, 4.85, 4.70, RGBColor(252, 252, 249))
    add_metric(s, f"{metrics['magnetic_roc_auc']:.3f}", "磁异常检测AUC", 1.10, 1.92, 2.6, TEAL)
    add_bullets(
        s,
        [
            "WMM2025给出背景磁场尺度，合成偶极异常模拟海缆局部扰动。",
            "热力图中的异常高值区域被映射为风险优先级，进入任务权重。",
            "AUC接近1说明异常区域可区分性强，后续调度和轨迹优化因此具备明确的风险输入。",
        ],
        1.12,
        2.85,
        4.25,
        2.55,
        size=12.8,
        accent=CORAL,
    )
    add_picture(s, FIG / "magnetic_heatmap.png", 6.20, 1.30, 5.90, 5.55)
    add_footer(s, page, text="Magnetic anomaly field and priority mapping")
    page += 1

    s = new_slide()
    add_title(s, "RESULT 15", "鲁棒性、消融实验与Pareto权衡")
    add_picture(s, FIG / "monte_carlo_violin.png", 0.48, 1.42, 3.95, 2.75)
    add_picture(s, FIG / "ablation_bar.png", 4.68, 1.42, 3.95, 2.75)
    add_picture(s, FIG / "pareto_energy_secure_data.png", 8.88, 1.42, 3.95, 2.75)
    add_explanation_panel(
        s,
        "鲁棒性图",
        "横轴为四类算法，纵轴为安全能效。每个小提琴图来自80次随机信道扰动，宽度表示结果分布密度。Proposed分布整体最高且离散程度可控，说明联合优化对信道扰动具有稳定优势。",
        0.55,
        4.50,
        3.80,
        1.55,
        TEAL,
    )
    add_explanation_panel(
        s,
        "消融实验图",
        "横向条形图逐一移除QKD优先级、RIS优化、UAV轨迹和磁异常感知模块。Full模型最高；移除RIS或QKD后下降明显，说明传播增强和密钥调度是安全能效提升的核心来源。",
        4.75,
        4.50,
        3.80,
        1.55,
        CORAL,
    )
    add_explanation_panel(
        s,
        "Pareto曲线图",
        "横轴为总能耗，纵轴为加权安全数据量，颜色表示能耗价格参数。曲线前段增益较快，后段趋于平缓，说明继续增加能耗会面临边际收益递减。",
        8.95,
        4.50,
        3.80,
        1.55,
        BLUE,
    )
    add_footer(s, page, text="Robustness, ablation and Pareto analysis")
    page += 1

    image_slide(prs, page, "RESULT 16", "QKD供给与RIS规模的二维敏感性分析", FIG / "sensitivity_qkd_ris_heatmap.png", "该热力图同时改变QKD密钥供给和RIS单元数，观察最终安全能效。结果显示两个维度均能提升系统性能，但高密钥供给或大规模RIS区域出现边际收益递减。该现象说明系统优化需要同时考虑安全资源和传播资源，单独增加某一种资源会受另一类瓶颈限制。")
    page += 1

    image_slide(prs, page, "RESULT 17", "优化约束违反量检查", FIG / "constraint_violation_check.png", "该图以对数坐标展示速度、QKD密钥、功率和安全数据约束的最大违反量。所有违反量均低于预设容差，说明前述性能提升满足可行性要求。该检查是优化实验中必要的可行性验证。")
    page += 1

    transition_slide(
        prs,
        page,
        "六",
        "创新点与工程边界",
        "最后回到系统层面，总结本项目的创新闭环、应用路径和仍可继续推进的工程方向。",
        ["创新来自量子感知、量子安全和可重构传播的联合建模。", "工程部署需要把系统级仿真和单元级全波仿真继续打通。", "结果已经形成可复现代码、数据、图表和PPT闭环。"],
    )
    page += 1

    analysis_slide(
        prs,
        page,
        "DISCUSSION",
        "综合讨论：创新闭环与工程边界",
        [
            "前面的模型和实验说明，本项目的主要价值在于把三类通常分开讨论的技术放进同一个优化闭环：量子磁力计提供风险权重，QKD提供安全数据的密钥约束，RIS-UAV协同提供可控传播和移动中继能力。",
            "从工程角度看，系统输出是一组可执行控制量：UAV应飞向哪里、RIS应采用何种相位码本、哪些节点在何时被服务、功率和密钥如何分配。这些输出可以直接服务于灾后巡检和关键告警回传。",
            "同时，系统级仿真仍有边界：真实海况、全波单元结构、硬件非理想相位响应和多UAV协同都需要更细粒度建模。结尾部分将把贡献、局限和后续扩展分开说明。",
        ],
        "本章作用：把算法结果回收到应用部署和后续研究问题上。",
        dark=False,
    )
    page += 1

    s = new_slide()
    add_title(s, "APPLICATION", "面向灾后应急监测的部署流程")
    steps = [
        ("灾后态势汇聚", "卫星遥感与海上节点报告形成初始任务区域。"),
        ("异常优先排序", "AUV/浮标磁探测定位疑似海缆异常点。"),
        ("联合优化调度", "岸基计算中心求解UAV轨迹、RIS相位、资源和密钥分配。"),
        ("安全回传处置", "关键告警优先加密回传，电网中心生成巡检和修复决策。"),
    ]
    for i, (t, b) in enumerate(steps):
        x = 0.85 + i * 3.05
        add_panel(s, x, 1.75, 2.45, 3.35, RGBColor(252, 252, 249))
        add_text(s, f"{i + 1}", x + 0.20, 2.00, 0.40, 0.32, 17, TEAL, True, font=FONT_EN)
        add_text(s, t, x + 0.62, 1.98, 1.55, 0.34, 14.2, INK, True)
        add_text(s, b, x + 0.26, 2.55, 1.95, 1.42, 11.6, MUTED, leading=1.10)
        details = [
            "输入：灾害区域、风电节点和海况信息。",
            "输入：磁扰动、风险等级和数据新鲜度。",
            "输出：𝒒、𝜣、𝒑、𝒃、𝒂、𝒔、κ。",
            "输出：加密告警、巡检路径和修复优先级。",
        ]
        add_text(s, details[i], x + 0.26, 4.12, 1.95, 0.56, 10.8, INK, True, leading=1.08)
    add_text(s, "应用解释：该流程把遥感发现、磁异常定位、优化调度和安全回传串联起来。其关键意义在于，当远海链路、UAV电量和QKD密钥同时受限时，系统能够用优化模型明确决定哪些节点、哪些时隙、哪些数据应被优先服务，并把数学输出转化为可执行的巡检和通信控制指令。", 0.95, 5.55, 10.9, 0.88, 12.7, INK, True, leading=1.12)
    add_footer(s, page)
    page += 1

    # conclusion page
    s = new_slide(DEEP)
    add_title(s, "INNOVATION", "研究贡献、局限与后续扩展", dark=True)
    innovations = [
        ("贡献一：量子安全约束", "将QKD密钥预算显式写入安全数据分配问题。"),
        ("贡献二：量子磁感知权重", "将海缆磁异常结果转化为风险感知任务权重。"),
        ("贡献三：RIS-UAV联合优化", "联合轨迹、相位和资源分配以降低单位安全数据能耗。"),
        ("局限与扩展方向", "后续可接入CST/HFSS全波单元仿真和真实海况时序数据。"),
    ]
    for i, (t, b) in enumerate(innovations):
        x = 0.85 + (i % 2) * 5.95
        y = 1.70 + (i // 2) * 1.55
        add_panel(s, x, y, 5.10, 1.05, RGBColor(42, 58, 73), RGBColor(78, 95, 108))
        add_text(s, t, x + 0.25, y + 0.13, 4.55, 0.38, 14.2, GOLD if i == 0 else TEAL, True)
        add_text(s, b, x + 0.25, y + 0.58, 4.55, 0.34, 12.3, RGBColor(218, 226, 231))
    add_text(s, "总结：本项目将远海灾害监测问题转化为一个包含量子感知、量子安全、RIS可重构传播和UAV轨迹控制的联合优化模型，并通过可复现实验验证其链路收益、收敛性、约束可行性和鲁棒性。", 0.95, 5.38, 11.10, 0.72, 15.2, WHITE, True, leading=1.12)
    add_footer(s, page, text="Final innovation summary")
    page += 1

    total = len(prs.slides)
    if total != TOTAL_SLIDES:
        raise RuntimeError(f"Deck has {total} slides; expected {TOTAL_SLIDES}. Update TOTAL_SLIDES or page flow.")
    prs.save(FINAL)

    note_lines = ["# Speaker Notes V2", ""]
    for i, note in enumerate(notes, start=1):
        note_lines.extend([f"## Slide {i:02d}", note, ""])
    (REPORTS / "speaker_notes_v2.md").write_text("\n".join(note_lines), encoding="utf-8")
    print(f"Saved {FINAL} with {total} slides")
    return FINAL


if __name__ == "__main__":
    build()
